from dotenv import load_dotenv
load_dotenv()

import os
import re
import pyodbc
import pandas as pd
from typing import Annotated, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from openai import OpenAI
from typing import Tuple  # add this import if not already present

# === ENVIRONMENT ===
openai_api_key = os.environ.get("OPENAI_API_KEY")
# It's good practice to add a check for missing critical environment variables
if openai_api_key is None:
    raise ValueError("OPENAI_API_KEY environment variable not set.")
client = OpenAI(api_key=openai_api_key)
font_size = Pt(18)

# === SQL CONNECTION ===

def get_connection():
    return pyodbc.connect(
        f"Driver={{ODBC Driver 17 for SQL Server}};"
        f"Server={os.environ.get('SQL_SERVER')};"
        f"Database={os.environ.get('SQL_DATABASE')};"
        f"UID={os.environ.get('SQL_UID')};"
        f"PWD={os.environ.get('SQL_PASSWORD')};"
        "Encrypt=yes;TrustServerCertificate=yes;"
    )
# === GPT HELPERS ===
def generate_summary_from_dataframe(df: pd.DataFrame) -> str:
    prompt = f"""Summarize this dataset in 5 business bullet points (max 20 words each):\n{df.head(10).to_string(index=False)}"""
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a business analyst who writes summaries."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content.strip()

def generate_details_from_dataframe(df: pd.DataFrame) -> str:
    prompt = f"""List 3–5 specific facts or observations from this data:\n{df.head(10).to_string(index=False)}"""
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a data analyst."},
            {"role": "user", "content": prompt}
        ]
    )
    return response.choices[0].message.content.strip()

def generate_analysis_from_dataframe(df: pd.DataFrame) -> Tuple[str, str]:
    prompt = f"""Give 3–5 business insights and 2–3 recommendations from this dataset:\n{df.head(10).to_string(index=False)}"""
    response = client.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "system", "content": "You are a business insight generator."},
            {"role": "user", "content": prompt}
        ]
    )
    text = response.choices[0].message.content.strip()
    parts = re.split(r"\n{2,}", text)
    insights, recs = "", ""
    for part in parts:
        if "recommendation" in part.lower():
            recs = part
        else:
            insights = part
    return insights.strip(), recs.strip()

# === PPT GENERATION ===
def create_ppt(prompt: str, df: pd.DataFrame, summary: str, details: str, insights: str, recommendations: str) -> str:
    prs = Presentation()
    title_snippet = re.sub(r"[^\w\s]", "", prompt[:40]).strip().replace(" ", "")
    file_path = os.path.join(os.getcwd(), f"SupplySense_{title_snippet}.pptx")

    def add_bullet_slide(title, content):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for line in content.strip().split("\n")[:6]:
            if line.strip():
                p = tf.add_paragraph()
                p.text = f"• {line.strip()}"
                p.font.size = font_size
                p.level = 0

    def add_table_slide(title, data: pd.DataFrame):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        title_box.text_frame.text = title
        preview = data.head(7)
        rows, cols = preview.shape
        table = slide.shapes.add_table(rows + 1, cols, Inches(0.5), Inches(1), Inches(9), Inches(4.5)).table

        for i, col in enumerate(preview.columns):
            table.cell(0, i).text = str(col)
            table.cell(0, i).text_frame.paragraphs[0].font.size = font_size

        for r, row in enumerate(preview.itertuples(index=False), 1):
            for c, val in enumerate(row):
                table.cell(r, c).text = str(val)
                table.cell(r, c).text_frame.paragraphs[0].font.size = font_size
                table.cell(r, c).text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Slide 1: Title
    prs.slides.add_slide(prs.slide_layouts[0]).shapes.title.text = "Supply Sense AI"

    # Slides 2–7
    add_bullet_slide("Data Summary", summary)
    add_bullet_slide("Data Details", details)
    add_table_slide("Data Preview", df)
    add_bullet_slide("Data Insights", insights)
    add_bullet_slide("Recommendations", recommendations)
    add_bullet_slide("Thank You", "Thankyou")

    prs.save(file_path)
    return file_path

# === MAIN TOOL: Used by Agent ===
def execute_sql(reflection: Annotated[str, "Why you wrote this SQL"],
                sql: Annotated[str, "SQL SELECT query"]) -> Dict:
    try:
        with get_connection() as conn:
            df = pd.read_sql(sql, conn)
            summary = generate_summary_from_dataframe(df)
            details = generate_details_from_dataframe(df)
            insights, recommendations = generate_analysis_from_dataframe(df)
            ppt_path = create_ppt(reflection, df, summary, details, insights, recommendations)
            return {
                "result": df.head(10).to_dict(orient="records"),
                "ppt": ppt_path
            }
    except Exception as e:
        return {"error": str(e)}

